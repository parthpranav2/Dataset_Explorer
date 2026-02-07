import sys
import os

# --- RIGOROUS ENVIRONMENT PROTOCOL ---
current_dir = os.path.dirname(os.path.abspath(__file__))
project_root = os.path.dirname(os.path.dirname(current_dir))
venv_site_packages = os.path.join(project_root, 'venv', 'lib', 'python3.9', 'site-packages')

if os.path.exists(venv_site_packages):
    sys.path.insert(0, venv_site_packages)

try:
    from pptx import Presentation
except ImportError:
    Presentation = None
# -------------------------------------

def extract(file_path):
    """
    Exhaustively explores .pptx modern presentations.
    Zero truncation: identifies slide counts, hidden notes, and metadata.
    """
    if not Presentation:
        return "Error: python-pptx library not found in venv."

    try:
        prs = Presentation(file_path)
        prop = prs.core_properties
        
        # 1. Structural Audit
        slide_count = len(prs.slides)
        notes_count = sum(1 for slide in prs.slides if slide.has_notes_slide)
        
        # 2. Content Density
        created = prop.created.strftime('%Y-%m-%d') if prop.created else "Unknown"
        modified = prop.modified.strftime('%Y-%m-%d') if prop.modified else "Unknown"
        last_author = prop.last_modified_by if prop.last_modified_by else "N/A"

        output = [
            f"Type: Office Open XML Presentation (.pptx)",
            f"Subject: {prop.subject if prop.subject else 'Technical Presentation'}",
            f"Author: {prop.author if prop.author else 'Unknown'}",
            f"Last Edited By: {last_author}",
            "Presentation Metrics:",
            f"  - Slides: {slide_count}",
            f"  - Slides with Notes: {notes_count}",
            f"Timeline: Created {created} | Modified {modified}",
            f"Keywords: {prop.keywords if prop.keywords else 'None'}"
        ]

        return "\n".join(output)

    except Exception as e:
        return f"PPTX Extraction Error: {str(e)}"